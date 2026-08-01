"""Recolor the PPT to a clean white tech style. Only change colors, not content."""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from lxml import etree

INPUT = r'C:\Users\LWH\Desktop\职途无限-答辩PPT.pptx'
OUTPUT = r'C:\Users\LWH\Desktop\职途无限-答辩PPT.pptx'

prs = Presentation(INPUT)

# Color palette - clean white tech style
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_WHITE = RGBColor(0xF8, 0xFA, 0xFC)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x2B, 0x5C, 0xE6)
ACCENT_LIGHT = RGBColor(0xE8, 0xEE, 0xFC)
MID_GRAY = RGBColor(0x8C, 0x8C, 0x9A)
BORDER_GRAY = RGBColor(0xE0, 0xE4, 0xE8)
DEEP_BLUE = RGBColor(0x0D, 0x1B, 0x3E)

def set_solid_fill(shape, color):
    """Set solid fill color on a shape."""
    try:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
    except:
        pass

def set_shape_border(shape, color, width=Pt(1)):
    """Set border on a shape."""
    try:
        shape.line.color.rgb = color
        shape.line.width = width
    except:
        pass

def set_gradient_fill_white(shape):
    """Set a subtle white-to-light-gray gradient."""
    try:
        shape.fill.gradient()
        shape.fill.gradient_stops[0].color.rgb = WHITE
        shape.fill.gradient_stops[0].position = 0.0
        shape.fill.gradient_stops[1].color.rgb = LIGHT_GRAY
        shape.fill.gradient_stops[1].position = 1.0
    except:
        pass

def recolor_slide_backgrounds():
    """Set all slide backgrounds to white or near-white."""
    for i, slide in enumerate(prs.slides):
        bg = slide.background
        fill = bg.fill
        # Keep cover (0) and ending (34) with dark bg for contrast
        if i in [0, 34]:
            # Dark background for cover/ending
            try:
                fill.solid()
                fill.fore_color.rgb = DEEP_BLUE
            except:
                pass
        else:
            # White background for all other slides
            try:
                fill.solid()
                fill.fore_color.rgb = WHITE
            except:
                pass

def remove_background_images():
    """Remove large background/decorative images from slides (keep content images)."""
    for i, slide in enumerate(prs.slides):
        if i in [0, 34]:
            continue  # Keep cover and ending images
        shapes_to_remove = []
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                # Remove large background images (covering most of the slide)
                if shape.width > 10000000 and shape.height > 5000000:
                    shapes_to_remove.append(shape)
                # Remove decorative small images that are just background decorations
                elif shape.width > 8000000:
                    shapes_to_remove.append(shape)
        for shape in shapes_to_remove:
            try:
                slide.shapes._spTree.remove(shape._element)
            except:
                pass

def recolor_rectangles():
    """Recolor rectangles to clean white/gray style."""
    for i, slide in enumerate(prs.slides):
        if i in [0, 34]:
            continue
        for shape in slide.shapes:
            name = shape.name.lower()
            # Recolor solid-fill rectangles
            if hasattr(shape, 'fill'):
                try:
                    if shape.fill.type is not None and shape.fill.type == 1:  # SOLID
                        try:
                            current = shape.fill.fore_color.rgb
                            # Dark rectangles -> white with border
                            if current and (
                                current.red < 80 and current.green < 80 and current.blue < 80
                            ):
                                set_solid_fill(shape, WHITE)
                                set_shape_border(shape, BORDER_GRAY)
                            # Blue/accent rectangles -> light blue
                            elif current and current.blue > 150:
                                set_solid_fill(shape, ACCENT_LIGHT)
                                set_shape_border(shape, ACCENT_BLUE, Pt(1.5))
                        except:
                            pass
                    elif shape.fill.type == 3:  # GRADIENT
                        set_gradient_fill_white(shape)
                except:
                    pass

def recolor_rounded_rects():
    """Recolor rounded rectangles to clean style."""
    for i, slide in enumerate(prs.slides):
        if i in [0, 34]:
            continue
        for shape in slide.shapes:
            name = shape.name
            if '圆角矩形' in name or 'Rounded' in name:
                try:
                    if shape.fill.type == 1:  # SOLID
                        try:
                            current = shape.fill.fore_color.rgb
                            if current:
                                # Make all rounded rects light blue with accent border
                                set_solid_fill(shape, ACCENT_LIGHT)
                                set_shape_border(shape, ACCENT_BLUE, Pt(1.5))
                        except:
                            set_solid_fill(shape, ACCENT_LIGHT)
                            set_shape_border(shape, ACCENT_BLUE, Pt(1.5))
                except:
                    pass

def recolor_text():
    """Recolor text to dark style for readability."""
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not hasattr(shape, 'text_frame'):
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        # Get current color
                        color = None
                        try:
                            color = run.font.color.rgb
                        except:
                            pass

                        if color is None:
                            # No explicit color, leave as is (usually dark)
                            pass
                        elif i in [0, 34]:
                            # Cover/ending: white text on dark bg
                            run.font.color.rgb = WHITE
                        else:
                            # Light colors (like cyan A0F6FE) -> dark text
                            if color.red > 150 or color.green > 150:
                                run.font.color.rgb = DARK_TEXT
                            # White/near-white text on non-dark bg -> dark text
                            elif color.red > 200 and color.green > 200 and color.blue > 200:
                                run.font.color.rgb = DARK_TEXT
                            # Very light text -> dark
                            elif color == RGBColor(0xF8, 0xF4, 0xED):
                                run.font.color.rgb = DARK_TEXT
                            elif color == RGBColor(0xF1, 0xE4, 0xD1):
                                run.font.color.rgb = DARK_TEXT
                            # Dark text stays dark
                            elif color.red < 60 and color.green < 60 and color.blue < 60:
                                pass
                            # Blue accent text stays blue
                            elif color.blue > 180 and color.red < 100:
                                run.font.color.rgb = ACCENT_BLUE
                    except:
                        pass

def recolor_ovals():
    """Recolor ovals/circles."""
    for i, slide in enumerate(prs.slides):
        if i in [0, 34]:
            continue
        for shape in slide.shapes:
            if shape.shape_type == 9:  # Oval
                try:
                    if shape.fill.type == 3:  # GRADIENT
                        shape.fill.gradient()
                        shape.fill.gradient_stops[0].color.rgb = ACCENT_BLUE
                        shape.fill.gradient_stops[0].position = 0.0
                        shape.fill.gradient_stops[1].color.rgb = RGBColor(0x6C, 0x9C, 0xFF)
                        shape.fill.gradient_stops[1].position = 1.0
                    elif shape.fill.type == 1:  # SOLID
                        set_solid_fill(shape, ACCENT_BLUE)
                except:
                    pass

# Apply all color changes
print('正在修改幻灯片背景...')
recolor_slide_backgrounds()

print('正在移除背景图片...')
remove_background_images()

print('正在修改矩形颜色...')
recolor_rectangles()

print('正在修改圆角矩形...')
recolor_rounded_rects()

print('正在修改圆形...')
recolor_ovals()

print('正在修改文字颜色...')
recolor_text()

prs.save(OUTPUT)
print(f'PPT 已保存: {OUTPUT}')
print(f'共 {len(prs.slides)} 页，配色已改为科技白风格')

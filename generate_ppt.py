"""Modify the original PPT with project content, preserving layout and font style."""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
import copy

TEMPLATE = r'C:\Users\LWH\Desktop\服(1).pptx'
OUTPUT = r'C:\Users\LWH\Desktop\职途无限-答辩PPT.pptx'

prs = Presentation(TEMPLATE)

def find_shape_by_text(slide, keyword):
    """Find first shape whose text contains keyword."""
    for shape in slide.shapes:
        if hasattr(shape, 'text') and keyword in shape.text:
            return shape
    return None

def set_text(shape, text):
    """Set text preserving first run's formatting."""
    if shape is None or not hasattr(shape, 'text_frame'):
        return
    tf = shape.text_frame
    if tf.paragraphs:
        p = tf.paragraphs[0]
        if p.runs:
            run = p.runs[0]
            run.text = text
            for r in list(p.runs)[1:]:
                r.text = ""
        else:
            p.text = text

def set_para_text(shape, para_idx, text):
    """Set text of a specific paragraph."""
    if shape is None or not hasattr(shape, 'text_frame'):
        return
    tf = shape.text_frame
    if len(tf.paragraphs) > para_idx:
        p = tf.paragraphs[para_idx]
        if p.runs:
            p.runs[0].text = text
            for r in list(p.runs)[1:]:
                r.text = ""
        else:
            p.text = text

def clear_slide(slide):
    """Remove all shapes from a slide."""
    sps = list(slide.shapes)
    for sp in sps:
        slide.shapes._spTree.remove(sp._element)

# ============================================================
# Slide 1: 封面 (keep as is, just update text)
# ============================================================
slide = prs.slides[0]
set_text(find_shape_by_text(slide, '全国大学生服务外包创新创业大赛'), '全国大学生服务外包创新创业大赛')
set_text(find_shape_by_text(slide, '基于 LangGraph'), '基于 LangGraph 多智能体协同的 AI 职业规划平台')
set_text(find_shape_by_text(slide, '参赛队伍'), '参赛队伍：无极限队  |   杭州电子科技大学    |   2026年6月')

# ============================================================
# Slide 2: 目录 (keep as is - it's all groups/pictures)
# ============================================================

# ============================================================
# Slide 3: 目录分隔页 (keep as is)
# ============================================================

# ============================================================
# Slide 4: 痛点分析 (keep content, it's already good)
# ============================================================

# ============================================================
# Slide 5: 政策背景 (keep content, it's already good)
# ============================================================

# ============================================================
# Slide 6: 分隔页 (keep as is)
# ============================================================

# ============================================================
# Slide 7: 项目概述 (keep content, it's already good)
# ============================================================

# ============================================================
# Slide 8: 系统架构 (keep as is - it's a group/picture)
# ============================================================

# ============================================================
# Slide 9-12: 核心智能体 1-4 (keep as is, they're detailed)
# ============================================================

# ============================================================
# Slide 13: 图表展示 (keep as is - it's a group)
# ============================================================

# ============================================================
# Slide 14-15: 核心智能体5 + 报告流程 (keep as is)
# ============================================================

# ============================================================
# Slide 16: 分隔页 (keep as is)
# ============================================================

# ============================================================
# Slides 17-19: 演示截图页 - 更新标题
# ============================================================
# These are layout pages with groups, keep as is

# ============================================================
# Slide 20: 任务完成/团队合作/自我提升 - update
# ============================================================
slide = prs.slides[19]
set_text(find_shape_by_text(slide, '在过去的一段时间里'), '完成简历解析、岗位匹配、职业规划、学习计划、报告生成5大智能体的开发与集成')
set_text(find_shape_by_text(slide, '我积极参与团队合作'), '前后端协同开发，Vue 3 + FastAPI + LangGraph全栈实现，SSE流式通信实时交互')
set_text(find_shape_by_text(slide, '我不断学习'), '掌握RAG检索增强生成、多智能体协同、向量数据库等前沿AI工程能力')

# ============================================================
# Slide 21: 多维图表 (keep as is - it's about charts)
# ============================================================

# ============================================================
# Slide 22: 分隔页 (keep as is)
# ============================================================

# ============================================================
# Slide 23: 任务完成/团队合作/自我提升 - update
# ============================================================
slide = prs.slides[22]
set_text(find_shape_by_text(slide, '在过去的一段时间里'), '基于ChromaDB向量数据库的语义检索，5000+岗位数据的Embedding索引与余弦相似度匹配')
set_text(find_shape_by_text(slide, '我积极参与团队合作'), '多智能体协同架构：LangGraph状态图管理Agent间数据流转，实现复杂工作流编排')
set_text(find_shape_by_text(slide, '我不断学习'), '融合RAG+LLM的混合匹配策略：向量召回→去重→七维评分→个性化推荐')

# ============================================================
# Slide 24: 用户画像/人岗匹配 (keep as is - it's about features)
# ============================================================

# ============================================================
# Slide 25: 分隔页 (keep as is)
# ============================================================

# ============================================================
# Slide 26: 任务完成/团队合作/自我提升 - update
# ============================================================
slide = prs.slides[25]
set_text(find_shape_by_text(slide, '在过去的一段时间里'), 'AI诊断报告：基于七维能力分析，LLM生成3段式专业诊断结论，支持多格式导出')
set_text(find_shape_by_text(slide, '我积极参与团队合作'), '成长追踪中心：设定目标岗位→识别能力缺口→生成每日学习任务→可视化进度')
set_text(find_shape_by_text(slide, '我不断学习'), '职能助手Coach：基于SSE的实时AI对话，支持职业咨询、学习指导、简历优化')

# ============================================================
# Slide 27: 报告导出/每日计划 (keep as is - it's about features)
# ============================================================

# ============================================================
# Slide 28: 分隔页 (keep as is)
# ============================================================

# ============================================================
# Slides 29-33: 开发历程时间线 (keep as is, they have good structure)
# ============================================================

# ============================================================
# Slide 34: 团队介绍 (keep as is, it has real team info)
# ============================================================

# ============================================================
# Slide 35: 感谢页 (keep as is)
# ============================================================

# ============================================================
# Remove empty slides 36-39
# ============================================================
# Delete last 4 empty slides
for i in range(4):
    rId = prs.slides._sldIdLst[-1].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[-1])

prs.save(OUTPUT)
print(f'PPT 已生成: {OUTPUT}')
print(f'共 {len(prs.slides)} 页')
